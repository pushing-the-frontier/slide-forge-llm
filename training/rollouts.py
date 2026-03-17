"""Generate multi-step trajectories using Claude via AWS Bedrock.

Reward System:
- Step rewards are computed using the actual reward functions (code_rules, content_quality, etc.)
- After each action that modifies slides, we compute quality_delta = new_quality - old_quality
- This provides dense, meaningful reward signals for RL training
"""

from __future__ import annotations

import argparse
import json
import os
import time
from pathlib import Path
from collections import defaultdict
from dotenv import load_dotenv

load_dotenv(Path(__file__).parent.parent / ".env")

import boto3

from envs.slideforge_env.models import SlideForgeAction, SlideBrief, SlideForgeState
from envs.slideforge_env.server.environment import SlideForgeEnvironment
from rewards.aggregator import compute_reward_details, DEFAULT_WEIGHTS
from training.grpo_trainer import extract_tool_call
from training.prompts import format_prompt

# Config
BEDROCK_MODEL_ID = "us.anthropic.claude-opus-4-6-v1"
BEDROCK_REGION = os.environ.get("AWS_REGION", "us-east-1")
MAX_TOKENS_PER_TURN = 1024

SYSTEM_PROMPT = """You are SlideForge, an expert AI agent that creates and edits professional business presentations.

You specialize in creating data-driven, executive-ready slide decks for:
- Investor pitch decks (Seed, Series A/B/C)
- Board meetings and quarterly business reviews
- Financial analysis and budget presentations
- Strategic planning and market analysis
- Product roadmaps and go-to-market strategies

## Content Guidelines
- Use SPECIFIC numbers, percentages, and metrics from research or provided data
- Never make up statistics - only use data from web_search results or provided content
- Include realistic timeframes and milestones
- Structure content for the specific audience
- Use professional business terminology appropriate to the topic

## Strategy: Creating New Presentations
1. If no source data is provided, use web_search 2-3 times to gather real statistics and facts
2. Create an outline with create_outline covering all target slides
3. Set an appropriate theme with set_theme (corporate, tech, dark, creative)
4. Generate each slide with generate_slide using researched/provided data
5. Review with review_deck and fix any issues with edit_slide
6. Call finalize when done

## Strategy: Editing Existing Presentations
1. Use get_slide_content on each slide to understand the current deck
2. Identify weaknesses: vague content, missing data, poor flow, structural issues
3. Use web_search to gather fresh data for slides that need stronger content
4. Make targeted edits with edit_slide (content), or structural changes with insert_slide, delete_slide, reorder_slides
5. Review with review_deck and finalize when done

## Rules
- Always respond with EXACTLY ONE JSON tool call per turn
- Wrap your tool call in ```json ... ``` markers
- Think briefly before each tool call (1-2 sentences), then output the JSON
- Generate ALL target slides - do not skip any
- Use only factual data from research or provided content"""

DEFAULT_BRIEFS = [
    {"topic": "Artificial Intelligence", "audience": "executives", "num_slides": 5, "sections_per_slide": 3},
]


# =============================================================================
# Bedrock Client
# =============================================================================

def _create_bedrock_client(region: str | None = None):
    return boto3.client("bedrock-runtime", region_name=region or BEDROCK_REGION)


def _call_claude(
    client,
    messages: list[dict],
    system: str,
    model_id: str = BEDROCK_MODEL_ID,
    max_tokens: int = MAX_TOKENS_PER_TURN,
    temperature: float = 0.7,
) -> str:
    """Call Claude via Bedrock Converse API."""
    response = client.converse(
        modelId=model_id,
        system=[{"text": system}],
        messages=messages,
        inferenceConfig={"maxTokens": max_tokens, "temperature": temperature},
    )
    output_message = response["output"]["message"]
    return "\n".join(block["text"] for block in output_message["content"] if "text" in block)


# =============================================================================
# Step Reward Computation (Using Actual Reward Functions)
# =============================================================================

def compute_quality_score(state: SlideForgeState, weights: dict | None = None) -> dict:
    """Compute current quality score using the reward functions."""
    w = weights or DEFAULT_WEIGHTS
    try:
        details = compute_reward_details(["x"], weights=w, states=[state])
        return details[0]
    except Exception:
        return {"aggregate": 0.0, "code_rules": 0.0, "render_quality": 0.0, 
                "content_quality": 0.0, "claude_aesthetic_html": 0.0,
                "claude_aesthetic_visual": 0.0,
                "brief_reconstruction": 0.0}


def compute_step_reward(
    state: SlideForgeState,
    prev_quality: dict,
    tool: str,
    success: bool,
    weights: dict | None = None,
) -> dict:
    """Compute step reward as quality delta plus action bonuses.
    
    Returns:
        dict with step_reward and component breakdown
    """
    # Compute current quality
    curr_quality = compute_quality_score(state, weights)
    
    # Quality delta (main reward signal)
    quality_delta = curr_quality["aggregate"] - prev_quality["aggregate"]
    
    # Small bonuses for successful actions (encourages progress)
    action_bonus = 0.0
    if success:
        action_bonus = 0.01  # Small bonus for any successful action
        if tool == "finalize" and state.phase == "DONE":
            action_bonus = 0.1  # Completion bonus
    else:
        action_bonus = -0.02  # Small penalty for failed actions
    
    # Phase progression bonus (encourages moving forward)
    phase_bonus = 0.0
    phase_values = {"RESEARCH": 0, "PLAN": 0.02, "GENERATE": 0.05, "REFINE": 0.08, "DONE": 0.1}
    phase_bonus = phase_values.get(state.phase, 0)
    
    step_reward = quality_delta + action_bonus
    
    return {
        "step_reward": step_reward,
        "quality_delta": quality_delta,
        "action_bonus": action_bonus,
        "phase_bonus": phase_bonus,
        "current_quality": curr_quality,
        "prev_quality_aggregate": prev_quality["aggregate"],
    }


# =============================================================================
# Rollout Execution
# =============================================================================

def run_rollout(
    brief: dict,
    client=None,
    model_id: str = BEDROCK_MODEL_ID,
    max_turns: int = 30,
    temperature: float = 0.7,
    reward_weights: dict | None = None,
    verbose: bool = False,
    existing_slides: list[str] | None = None,
) -> dict:
    """Run a single rollout episode.
    
    Args:
        existing_slides: If provided, starts in edit mode with these slides pre-loaded.
    """
    if client is None:
        client = _create_bedrock_client()

    env = SlideForgeEnvironment()
    obs = env.reset(brief=brief, existing_slides=existing_slides)
    
    content = brief.get("content") if isinstance(brief, dict) else None
    user_context = format_prompt(env.state, content=content)

    messages = [{"role": "user", "content": [{"text": user_context}]}]
    
    # Initialize quality tracking
    prev_quality = compute_quality_score(env.state, reward_weights)
    cumulative_reward = 0.0

    trajectory = {
        "episode_id": env.state.episode_id,
        "brief": brief,
        "turns": [],
        "total_steps": 0,
        "final_phase": "RESEARCH",
        "completed": False,
    }

    for turn_idx in range(max_turns):
        # Call Claude
        try:
            assistant_text = _call_claude(
                client, messages, SYSTEM_PROMPT, model_id, MAX_TOKENS_PER_TURN, temperature
            )
        except Exception as e:
            if verbose:
                print(f"  [turn {turn_idx}] Bedrock error: {e}")
            trajectory["turns"].append({
                "turn": turn_idx, "assistant": f"ERROR: {e}", "tool_call": None,
                "observation": None, "success": False,
                "step_reward": {"step_reward": -0.05, "error": str(e)},
            })
            break

        # Extract tool call
        tool_call = extract_tool_call(assistant_text)
        if tool_call is None:
            if verbose:
                print(f"  [turn {turn_idx}] No tool call extracted")
            messages.append({"role": "assistant", "content": [{"text": assistant_text}]})
            messages.append({"role": "user", "content": [{"text": "You must respond with a JSON tool call."}]})
            trajectory["turns"].append({
                "turn": turn_idx, "assistant": assistant_text, "tool_call": None,
                "observation": "No tool call extracted.", "success": False,
                "step_reward": {"step_reward": -0.01, "parse_error": True},
            })
            continue

        # Execute action -- handle LLMs that put params at top-level instead of under "parameters"
        params = tool_call.get("parameters", {})
        if not params:
            params = {k: v for k, v in tool_call.items() if k != "tool"}
        action = SlideForgeAction(tool=tool_call.get("tool", ""), parameters=params)
        obs = env.step(action)

        # Compute step reward using quality delta
        step_reward = compute_step_reward(
            state=env.state,
            prev_quality=prev_quality,
            tool=tool_call.get("tool", ""),
            success=obs.success,
            weights=reward_weights,
        )
        cumulative_reward += step_reward["step_reward"]
        prev_quality = step_reward["current_quality"]

        turn_record = {
            "turn": turn_idx,
            "assistant": assistant_text,
            "tool_call": tool_call,
            "observation": obs.result,
            "success": obs.success,
            "phase": obs.phase,
            "slide_count": obs.current_slide_count,
            "done": obs.done,
            "step_reward": step_reward,
            "cumulative_reward": cumulative_reward,
        }
        trajectory["turns"].append(turn_record)

        if verbose:
            status = "OK" if obs.success else "FAIL"
            q = step_reward["current_quality"]["aggregate"]
            print(f"  [turn {turn_idx}] {tool_call['tool']}: {status} | "
                  f"phase={obs.phase} slides={obs.current_slide_count} | "
                  f"r={step_reward['step_reward']:.3f} q={q:.3f}")

        # Update conversation
        messages.append({"role": "assistant", "content": [{"text": assistant_text}]})
        turns_left = max_turns - turn_idx - 1
        obs_text = (f"Tool result (success={obs.success}):\n{obs.result}\n\n"
                    f"State: phase={obs.phase}, slides={obs.current_slide_count}/{brief.get('num_slides', 10)}, "
                    f"turns remaining={turns_left}")
        if obs.done:
            obs_text += "\n\nEpisode complete."
        elif turns_left <= 3:
            obs_text += f"\n\nURGENT: Only {turns_left} turn(s) left. Call finalize NOW to complete the presentation."
        elif turns_left <= 6:
            obs_text += f"\n\nReminder: {turns_left} turns left. Wrap up edits and call finalize soon."
        messages.append({"role": "user", "content": [{"text": obs_text}]})

        if obs.done:
            break

    # Final trajectory data
    state = env.state
    trajectory["total_steps"] = state.step_count
    trajectory["final_phase"] = state.phase
    trajectory["completed"] = state.done and state.phase == "DONE"
    trajectory["slides_created"] = sum(1 for h in state.slides_html if h)
    trajectory["theme"] = state.theme
    trajectory["cumulative_reward"] = cumulative_reward
    trajectory["slides_html"] = state.slides_html
    trajectory["slides_png"] = state.slides_png
    trajectory["outline"] = state.outline
    trajectory["research_context"] = state.research_context
    trajectory["edit_mode"] = state.edit_mode
    trajectory["original_slides_html"] = state.original_slides_html if state.edit_mode else []
    trajectory["final_quality"] = compute_quality_score(state, reward_weights)
    trajectory["messages"] = _build_messages_list(trajectory)

    return trajectory


def _build_messages_list(trajectory: dict) -> list[dict]:
    """Convert trajectory to standard messages format for training."""
    messages = [{"role": "system", "content": SYSTEM_PROMPT}]
    brief = trajectory["brief"]
    brief_fields = {k: v for k, v in brief.items() if k != "content"}
    is_edit = trajectory.get("edit_mode", False)
    state = SlideForgeState(
        brief=SlideBrief(**brief_fields),
        edit_mode=is_edit,
        slides_html=list(trajectory.get("original_slides_html", [])) if is_edit else [],
    )
    content = brief.get("content") if isinstance(brief, dict) else None
    messages.append({"role": "user", "content": format_prompt(state, content=content)})

    for turn in trajectory["turns"]:
        if turn["assistant"]:
            messages.append({"role": "assistant", "content": turn["assistant"]})
        if turn["observation"]:
            messages.append({"role": "user", "content": f"Tool result (success={turn['success']}):\n{turn['observation']}"})
    return messages


# =============================================================================
# Batch Execution
# =============================================================================

def run_batch(
    briefs: list[dict] | None = None,
    num_rollouts_per_brief: int = 1,
    model_id: str = BEDROCK_MODEL_ID,
    max_turns: int = 30,
    temperature: float = 0.7,
    reward_weights: dict | None = None,
    verbose: bool = True,
    checkpoint_file: str | None = None,
    checkpoint_interval: int = 10,
    existing_slides_map: dict[int, list[str]] | None = None,
) -> list[dict]:
    """Run rollouts for a batch of briefs with checkpointing.
    
    Args:
        existing_slides_map: Optional mapping of brief_idx -> list of slide HTML strings.
            When present, those briefs run in edit mode with the given slides pre-loaded.
    """
    briefs = briefs or DEFAULT_BRIEFS
    client = _create_bedrock_client()
    all_trajectories = []
    
    # Load existing checkpoint if available
    completed_episodes = set()
    if checkpoint_file and os.path.exists(checkpoint_file):
        try:
            with open(checkpoint_file) as f:
                existing = json.load(f)
            # Only keep trajectories that actually completed; re-run incomplete ones
            all_trajectories = [t for t in existing if t.get("completed", False)]
            completed_episodes = {(t["brief"].get("topic"), t.get("rollout_idx", 0)) for t in all_trajectories}
            if verbose:
                n_skip = len(all_trajectories)
                n_retry = len(existing) - n_skip
                print(f"Loaded checkpoint: {n_skip} completed, {n_retry} incomplete (will retry)")
        except Exception as e:
            if verbose:
                print(f"Could not load checkpoint: {e}")

    total = len(briefs) * num_rollouts_per_brief
    skipped = 0
    
    for brief_idx, brief in enumerate(briefs):
        for rollout_idx in range(num_rollouts_per_brief):
            run_num = brief_idx * num_rollouts_per_brief + rollout_idx + 1
            
            # Skip if already completed
            key = (brief.get("topic"), rollout_idx)
            if key in completed_episodes:
                skipped += 1
                if verbose:
                    print(f"[{run_num}/{total}] Skipping (already done): {brief.get('topic', '?')}")
                continue
            
            if verbose:
                print(f"\n[{run_num}/{total}] Brief: {brief.get('topic', '?')} (rollout {rollout_idx + 1})")

            t0 = time.time()
            slides_for_brief = (existing_slides_map or {}).get(brief_idx)
            try:
                trajectory = run_rollout(
                    brief=brief, client=client, model_id=model_id,
                    max_turns=max_turns, temperature=temperature,
                    reward_weights=reward_weights, verbose=verbose,
                    existing_slides=slides_for_brief,
                )
            except Exception as e:
                if verbose:
                    print(f"  ERROR: {e}")
                trajectory = {
                    "episode_id": f"error_{brief_idx}_{rollout_idx}",
                    "brief": brief,
                    "turns": [],
                    "total_steps": 0,
                    "final_phase": "ERROR",
                    "completed": False,
                    "slides_created": 0,
                    "error": str(e),
                }
            
            elapsed = time.time() - t0
            trajectory["rollout_idx"] = rollout_idx
            trajectory["elapsed_seconds"] = round(elapsed, 1)
            all_trajectories.append(trajectory)

            if verbose:
                q = trajectory.get("final_quality", {}).get("aggregate", 0)
                status = "DONE" if trajectory.get("completed") else "INCOMPLETE"
                print(f"  {status} in {elapsed:.1f}s | slides={trajectory.get('slides_created', 0)} quality={q:.3f}")
            
            # Checkpoint periodically
            if checkpoint_file and (len(all_trajectories) - skipped) % checkpoint_interval == 0:
                _save_checkpoint(all_trajectories, checkpoint_file, verbose)
    
    # Final checkpoint
    if checkpoint_file:
        _save_checkpoint(all_trajectories, checkpoint_file, verbose)
    
    return all_trajectories


def _save_checkpoint(trajectories: list[dict], filepath: str, verbose: bool = True):
    """Save trajectories checkpoint."""
    os.makedirs(os.path.dirname(filepath) or ".", exist_ok=True)
    with open(filepath, "w") as f:
        json.dump(trajectories, f, indent=2, default=str)
    if verbose:
        print(f"  [checkpoint] Saved {len(trajectories)} trajectories to {filepath}")


# =============================================================================
# Dataset Conversion
# =============================================================================

def trajectories_to_dataset(trajectories: list[dict], format: str = "sft"):
    """Convert trajectories to HuggingFace Dataset."""
    from datasets import Dataset

    if format == "grpo":
        return _trajectories_to_grpo_dataset(trajectories)

    rows = []
    for traj in trajectories:
        fq = traj.get("final_quality", {})
        row = {
            "episode_id": traj["episode_id"],
            "topic": traj["brief"].get("topic", ""),
            "audience": traj["brief"].get("audience", ""),
            "num_slides_target": traj["brief"].get("num_slides", 10),
            "slides_created": traj["slides_created"],
            "completed": traj["completed"],
            "total_steps": traj["total_steps"],
            "final_phase": traj["final_phase"],
            "theme": traj.get("theme", "default"),
            "num_turns": len(traj["turns"]),
            "quality_aggregate": fq.get("aggregate", 0.0),
            "quality_code_rules": fq.get("code_rules", 0.0),
            "quality_render": fq.get("render_quality", 0.0),
            "quality_content": fq.get("content_quality", 0.0),
            "quality_aesthetic_html": fq.get("claude_aesthetic_html", 0.0),
            "quality_aesthetic_visual": fq.get("claude_aesthetic_visual", 0.0),
            "quality_brief_recon": fq.get("brief_reconstruction", 0.0),
            "cumulative_reward": traj.get("cumulative_reward", 0.0),
            "elapsed_seconds": traj.get("elapsed_seconds", 0.0),
        }
        if format == "sft":
            row["messages"] = traj["messages"]
        else:
            row["brief"] = json.dumps(traj["brief"])
            row["messages"] = json.dumps(traj["messages"])
            row["turns"] = json.dumps(traj["turns"])
        rows.append(row)

    return Dataset.from_list(rows)


def _trajectories_to_grpo_dataset(trajectories: list[dict]):
    """Convert to GRPO format with per-step data."""
    from datasets import Dataset

    all_steps = []
    for traj in trajectories:
        episode_id = traj["episode_id"]
        brief = traj["brief"]
        turns = traj["turns"]
        final_quality = traj.get("final_quality", {}).get("aggregate", 0.0)

        # Build conversation incrementally
        conversation = [{"role": "system", "content": SYSTEM_PROMPT}]
        brief_for_state = {k: v for k, v in brief.items() if k != "content"}
        state = SlideForgeState(brief=SlideBrief(**brief_for_state))
        content = brief.get("content") if isinstance(brief, dict) else None
        conversation.append({"role": "user", "content": format_prompt(state, content=content)})

        # Compute returns (reward-to-go)
        step_rewards = [t.get("step_reward", {}).get("step_reward", 0.0) for t in turns]
        returns = []
        G = final_quality
        for r in reversed(step_rewards):
            G = r + 0.99 * G
            returns.insert(0, G)

        for step_idx, turn in enumerate(turns):
            sr = turn.get("step_reward", {})
            step_data = {
                "episode_id": episode_id,
                "step_idx": step_idx,
                "topic": brief.get("topic", ""),
                "prompt": list(conversation),
                "completion": turn.get("assistant", ""),
                "tool": turn.get("tool_call", {}).get("tool", "") if turn.get("tool_call") else "",
                "success": turn.get("success", False),
                "phase": turn.get("phase", "RESEARCH"),
                "slide_count": turn.get("slide_count", 0),
                "step_reward": sr.get("step_reward", 0.0),
                "quality_delta": sr.get("quality_delta", 0.0),
                "current_quality": sr.get("current_quality", {}).get("aggregate", 0.0),
                "cumulative_reward": turn.get("cumulative_reward", 0.0),
                "return_to_go": returns[step_idx] if step_idx < len(returns) else 0.0,
                "episode_completed": traj["completed"],
                "episode_final_quality": final_quality,
            }
            all_steps.append(step_data)

            if turn.get("assistant"):
                conversation.append({"role": "assistant", "content": turn["assistant"]})
            if turn.get("observation"):
                conversation.append({"role": "user", "content": f"Tool result: {turn['observation']}"})

    # Compute advantages (GRPO-style grouping)
    step_groups = defaultdict(list)
    for i, step in enumerate(all_steps):
        step_groups[(step["topic"], step["step_idx"])].append(i)

    for step in all_steps:
        step["advantage"] = step["return_to_go"]

    for key, indices in step_groups.items():
        if len(indices) > 1:
            returns = [all_steps[i]["return_to_go"] for i in indices]
            mean_r, std_r = sum(returns) / len(returns), (sum((r - sum(returns)/len(returns))**2 for r in returns) / len(returns)) ** 0.5
            for i, idx in enumerate(indices):
                all_steps[idx]["advantage"] = (returns[i] - mean_r) / (std_r + 1e-8) if std_r > 1e-8 else 0.0

    return Dataset.from_list(all_steps)


def push_to_hub(trajectories: list[dict], repo_id: str, private: bool = True, token: str | None = None, format: str = "sft"):
    """Push trajectories to HuggingFace Hub."""
    dataset = trajectories_to_dataset(trajectories, format=format)
    hf_token = token or os.environ.get("HF_TOKEN")
    print(f"Pushing {len(dataset)} rows to {repo_id}...")
    dataset.push_to_hub(repo_id, private=private, token=hf_token)
    print(f"Dataset pushed to https://huggingface.co/datasets/{repo_id}")
    return dataset


# =============================================================================
# Slide Export
# =============================================================================

def _build_deck_html(trajectory: dict) -> str:
    """Build combined HTML deck using iframes."""
    topic = trajectory["brief"].get("topic", "Presentation")
    theme = trajectory.get("theme", "default")
    slides_html = trajectory.get("slides_html", [])

    slides_content = ""
    for i, html in enumerate(slides_html):
        if html:
            escaped = html.replace('"', '&quot;').replace("'", "&#39;")
            slides_content += f'''
        <div class="slide-container">
            <div class="slide-label">Slide {i + 1}</div>
            <iframe class="slide-frame" srcdoc="{escaped}" scrolling="no"></iframe>
        </div>'''

    return f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{topic}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: -apple-system, sans-serif; background: linear-gradient(135deg, #1a1a2e, #16213e); min-height: 100vh; padding: 40px 20px; }}
        .header {{ text-align: center; margin-bottom: 48px; }}
        h1 {{ color: #fff; font-size: 36px; margin-bottom: 12px; }}
        .meta {{ color: rgba(255,255,255,0.6); }}
        .meta span {{ display: inline-block; padding: 4px 12px; background: rgba(255,255,255,0.1); border-radius: 20px; margin: 0 4px; }}
        .slides-grid {{ max-width: 1400px; margin: 0 auto; display: flex; flex-direction: column; gap: 32px; }}
        .slide-container {{ background: rgba(255,255,255,0.05); border-radius: 16px; padding: 20px; border: 1px solid rgba(255,255,255,0.1); }}
        .slide-label {{ color: rgba(255,255,255,0.5); font-size: 12px; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 12px; }}
        .slide-frame {{ width: 100%; height: 540px; border: none; border-radius: 8px; background: #fff; box-shadow: 0 20px 60px rgba(0,0,0,0.3); }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{topic}</h1>
        <p class="meta"><span>Theme: {theme}</span><span>{len([h for h in slides_html if h])} Slides</span></p>
    </div>
    <div class="slides-grid">{slides_content}</div>
</body>
</html>'''


def _build_deck_pptx(trajectory: dict, output_path: str) -> None:
    """Build a .pptx file from a trajectory's slide data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from bs4 import BeautifulSoup

    from envs.slideforge_env.server.rendering.themes import THEMES

    topic = trajectory["brief"].get("topic", "Presentation")
    theme_name = trajectory.get("theme", "default")
    theme = THEMES.get(theme_name, THEMES["default"])
    slides_html = [h for h in trajectory.get("slides_html", []) if h]

    bg_rgb = RGBColor(*theme["bg"])
    text_rgb = RGBColor(*theme["text"])
    accent_rgb = RGBColor(*theme["accent"])
    secondary_rgb = RGBColor(*theme["secondary"])

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9 widescreen
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]  # blank layout

    for html in slides_html:
        soup = BeautifulSoup(html, "html.parser")

        title_el = soup.select_one(".title")
        title_text = title_el.get_text(strip=True) if title_el else ""

        sections = []
        for sec_el in soup.select(".section"):
            h2 = sec_el.select_one("h2")
            p = sec_el.select_one("p")
            sections.append({
                "heading": h2.get_text(strip=True) if h2 else "",
                "body": p.get_text(strip=True) if p else "",
            })

        slide = prs.slides.add_slide(blank_layout)

        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = bg_rgb

        # Title
        left, top = Inches(0.7), Inches(0.5)
        title_box = slide.shapes.add_textbox(left, top, Inches(8.6), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = accent_rgb

        # Accent bar under title
        slide.shapes.add_shape(
            1, Inches(0.7), Inches(1.35), Inches(1.0), Inches(0.06)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = accent_rgb
        slide.shapes[-1].line.fill.background()

        # Sections
        n = len(sections)
        if n == 0:
            continue

        content_top = Inches(1.7)
        content_height = Inches(4.5)
        margin = Inches(0.7)
        gap = Inches(0.25)
        total_width = Inches(8.6)

        if n == 1:
            col_width = total_width
        else:
            col_width = (total_width - gap * (n - 1)) / n

        for i, sec in enumerate(sections):
            col_left = margin + (col_width + gap) * i

            # Section heading
            heading_box = slide.shapes.add_textbox(col_left, content_top, col_width, Inches(0.5))
            htf = heading_box.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.text = sec["heading"]
            hp.font.size = Pt(16)
            hp.font.bold = True
            hp.font.color.rgb = accent_rgb

            # Section body
            body_box = slide.shapes.add_textbox(col_left, content_top + Inches(0.55), col_width, content_height - Inches(0.6))
            btf = body_box.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            bp.text = sec["body"]
            bp.font.size = Pt(13)
            bp.font.color.rgb = text_rgb
            bp.line_spacing = Pt(20)

    prs.save(output_path)


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(description="SlideForge Trajectory Rollouts")
    parser.add_argument("--briefs-file", help="JSON file with briefs")
    parser.add_argument("--topic", help="Single topic (overrides briefs-file)")
    parser.add_argument("--audience", default="general")
    parser.add_argument("--num-slides", type=int, default=5)
    parser.add_argument("--sections-per-slide", type=int, default=3)
    parser.add_argument("--num-rollouts", type=int, default=1)
    parser.add_argument("--limit", type=int, help="Limit number of briefs to process")
    parser.add_argument("--max-turns", type=int, default=30)
    parser.add_argument("--temperature", type=float, default=0.7)
    parser.add_argument("--model-id", default=BEDROCK_MODEL_ID)
    parser.add_argument("--region")
    parser.add_argument("--output", default="outputs/trajectories.json")
    parser.add_argument("--checkpoint", help="Checkpoint file for resumable runs (default: output file)")
    parser.add_argument("--checkpoint-interval", type=int, default=5, help="Save checkpoint every N rollouts")
    parser.add_argument("--save-slides", action="store_true")
    parser.add_argument("--push-to-hub")
    parser.add_argument("--hf-token")
    parser.add_argument("--dataset-format", default="sft", choices=["sft", "grpo", "raw"])
    parser.add_argument("--public", action="store_true")
    parser.add_argument("--quiet", action="store_true")
    parser.add_argument("--edit-from", help="Path to a previous trajectories JSON to load slides from for edit mode")
    parser.add_argument("--edit-instructions", default="", help="Instructions for the editing task")
    args = parser.parse_args()

    if args.region:
        global BEDROCK_REGION
        BEDROCK_REGION = args.region

    briefs = None
    if args.topic:
        briefs = [{"topic": args.topic, "audience": args.audience, 
                   "num_slides": args.num_slides, "sections_per_slide": args.sections_per_slide}]
    elif args.briefs_file:
        with open(args.briefs_file) as f:
            briefs = json.load(f)
        if args.limit:
            briefs = briefs[:args.limit]
        print(f"Loaded {len(briefs)} briefs from {args.briefs_file}")

    # Load existing slides for edit mode
    existing_slides_map = None
    if args.edit_from:
        with open(args.edit_from) as f:
            source_trajectories = json.load(f)
        existing_slides_map = {}
        for i, traj in enumerate(source_trajectories):
            slides = [h for h in traj.get("slides_html", []) if h]
            if slides:
                existing_slides_map[i] = slides
        # Build briefs from source trajectories if none provided
        if briefs is None:
            briefs = [t["brief"] for t in source_trajectories]
            if args.limit:
                briefs = briefs[:args.limit]
                existing_slides_map = {k: v for k, v in existing_slides_map.items() if k < args.limit}
        if args.edit_instructions:
            for b in briefs:
                b["edit_instructions"] = args.edit_instructions
        print(f"Edit mode: loaded slides from {len(existing_slides_map)} trajectories in {args.edit_from}")

    # Use output as checkpoint by default for resumability
    checkpoint_file = args.checkpoint or args.output

    trajectories = run_batch(
        briefs=briefs, num_rollouts_per_brief=args.num_rollouts,
        model_id=args.model_id, max_turns=args.max_turns,
        temperature=args.temperature, verbose=not args.quiet,
        checkpoint_file=checkpoint_file, checkpoint_interval=args.checkpoint_interval,
        existing_slides_map=existing_slides_map,
    )

    # Final save
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    with open(args.output, "w") as f:
        json.dump(trajectories, f, indent=2, default=str)
    print(f"\nSaved {len(trajectories)} trajectories to {args.output}")

    if args.save_slides:
        base_dir = "outputs/runs"
        os.makedirs(base_dir, exist_ok=True)
        
        for traj in trajectories:
            eid = traj["episode_id"][:8]
            topic = traj["brief"].get("topic", "unknown")
            # Create a clean folder name from topic
            slug = "".join(c if c.isalnum() or c in " -_" else "" for c in topic)
            slug = slug.replace(" ", "_").lower()[:50]
            folder_name = f"{slug}_{eid}"
            run_dir = f"{base_dir}/{folder_name}"
            os.makedirs(run_dir, exist_ok=True)
            
            # Save individual slides
            slides_dir = f"{run_dir}/slides"
            os.makedirs(slides_dir, exist_ok=True)
            for idx, html in enumerate(traj.get("slides_html", [])):
                if html:
                    with open(f"{slides_dir}/slide_{idx+1}.html", "w") as f:
                        f.write(html)
            
            # Save combined deck (HTML + PPTX)
            if traj.get("slides_html"):
                with open(f"{run_dir}/deck.html", "w") as f:
                    f.write(_build_deck_html(traj))
                try:
                    _build_deck_pptx(traj, f"{run_dir}/deck.pptx")
                except Exception as e:
                    print(f"  Warning: PPTX export failed: {e}")
            
            # Save trajectory for this brief
            traj_copy = {k: v for k, v in traj.items() if k not in ["slides_html", "slides_png"]}
            traj_copy["slides_count"] = len([h for h in traj.get("slides_html", []) if h])
            with open(f"{run_dir}/trajectory.json", "w") as f:
                json.dump(traj_copy, f, indent=2, default=str)
            
            # Save brief info
            with open(f"{run_dir}/brief.json", "w") as f:
                json.dump(traj["brief"], f, indent=2)
            
            print(f"  Saved: {run_dir}/")

    completed = sum(1 for t in trajectories if t["completed"])
    avg_quality = sum(t.get("final_quality", {}).get("aggregate", 0) for t in trajectories) / max(len(trajectories), 1)
    print(f"\nCompleted: {completed}/{len(trajectories)}")
    print(f"Avg quality: {avg_quality:.3f}")

    if args.push_to_hub:
        push_to_hub(trajectories, args.push_to_hub, private=not args.public, 
                    token=args.hf_token, format=args.dataset_format)


if __name__ == "__main__":
    main()
